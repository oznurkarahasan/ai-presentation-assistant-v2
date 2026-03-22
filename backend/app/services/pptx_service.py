"""
PPTX text extraction service with security validation.
"""
from pptx import Presentation
from fastapi import UploadFile
from app.core.exceptions import FileProcessingError, ValidationError
from app.core.logger import logger
from pypdf import PdfWriter, PdfReader
import re
import io
import os
import asyncio
import subprocess  # nosec B404
import shutil
import tempfile

# Security limits (same as PDF)
MAX_PPTX_SLIDES = 500
MAX_SLIDE_SIZE_KB = 5000  # 5MB per slide

def clean_text(text: str) -> str:
    """
    Cleans extracted PPTX text by removing null bytes and other invalid characters.
    """
    # Remove null bytes
    text = text.replace('\x00', '')
    # Remove other control characters except newlines and tabs
    text = re.sub(r'[\x01-\x08\x0b-\x0c\x0e-\x1f\x7f-\x9f]', '', text)
    # Normalize whitespace
    text = re.sub(r'\s+', ' ', text)
    return text.strip()

def validate_pptx_security(prs: Presentation, file_size: int) -> None:
    """
    Validates PPTX for security issues: slide bombs, excessive size.
    
    Raises:
        ValidationError: If PPTX fails security checks
    """
    # Check slide count (PPTX bomb protection)
    num_slides = len(prs.slides)
    if num_slides > MAX_PPTX_SLIDES:
        logger.warning(f"PPTX has too many slides: {num_slides}")
        raise ValidationError(
            f"PPTX has too many slides ({num_slides}). Maximum allowed: {MAX_PPTX_SLIDES}"
        )
    
    # Check average slide size (detect compression bombs)
    avg_slide_size = file_size / num_slides if num_slides > 0 else 0
    if avg_slide_size > MAX_SLIDE_SIZE_KB * 1024:
        logger.warning(f"PPTX has suspicious slide size: {avg_slide_size/1024:.2f}KB per slide")
        raise ValidationError(
            "PPTX file has unusually large slides. This may be a malicious file."
        )
    
    logger.debug(f"PPTX security validation passed: {num_slides} slides, {file_size/1024:.2f}KB")

async def extract_text_from_pptx(file: UploadFile, file_size: int = 0) -> tuple[list[str], str, float]:
    """
    Reads the PPTX and returns slide text with layout metadata.
    Extracts both slide text and speaker notes.
    
    Args:
        file: Uploaded PPTX file
        file_size: File size in bytes (for security validation)
        
    Returns:
        tuple[list[str], str, float]:
            - Slide texts extracted per slide
            - Orientation ('portrait' or 'landscape')
            - Aspect ratio (width / height)
    """
    try:
        # Read file content into memory
        file_content = await file.read()
        file_bytes = io.BytesIO(file_content)
        
        # Load presentation
        prs = Presentation(file_bytes)
        
        if len(prs.slides) == 0:
            raise FileProcessingError("PPTX file has no slides")
        
        # Security validation
        validate_pptx_security(prs, file_size)
        
        slides_text = []
        
        for i, slide in enumerate(prs.slides, 1):
            try:
                text_parts = []
                
                # Extract text from all shapes in slide
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        text_parts.append(shape.text)
                
                # Extract speaker notes if available
                if slide.has_notes_slide:
                    notes_frame = slide.notes_slide.notes_text_frame
                    if notes_frame and notes_frame.text:
                        text_parts.append(f"Notes: {notes_frame.text}")
                
                # Combine all text from slide
                slide_text = "\n".join(text_parts) if text_parts else ""
                cleaned_text = clean_text(slide_text)
                slides_text.append(cleaned_text)
                
                logger.debug(f"Extracted slide {i}/{len(prs.slides)}")
                
            except Exception as e:
                logger.warning(f"Failed to extract slide {i}: {str(e)}")
                slides_text.append("")  # Add empty string for failed slides
        
        logger.info(f"Successfully extracted {len(slides_text)} slides from PPTX")
        
        # Get orientation
        width = prs.slide_width
        height = prs.slide_height
        orientation = "portrait" if height > width else "landscape"
        
        aspect_ratio = width / height if height > 0 else 1.777
        return slides_text, orientation, aspect_ratio

    except Exception as e:
        logger.error(f"PPTX extraction error: {str(e)}", exc_info=True)
        raise FileProcessingError(
            message="Failed to extract text from PPTX",
            details=str(e)
        )

def _strip_pdf_bookmarks(pdf_path: str) -> None:
    """Remove outline/bookmarks from PDF so the browser navpane doesn't auto-open."""
    try:
        reader = PdfReader(pdf_path)
        writer = PdfWriter()
        for page in reader.pages:
            writer.add_page(page)
        with open(pdf_path, "wb") as f:
            writer.write(f)
    except Exception as e:
        logger.warning(f"Failed to strip PDF bookmarks from {pdf_path}: {e}")


async def convert_to_pdf_preview(pptx_path: str) -> str | None:
    """
    Converts a PPTX file to a PDF preview using LibreOffice headless.
    Output is saved as {pptx_path}.preview.pdf next to the original file.
    Returns the preview path on success, None on failure (non-breaking).
    """
    abs_pptx = os.path.abspath(pptx_path)
    abs_outdir = os.path.dirname(abs_pptx)
    preview_path = abs_pptx + ".preview.pdf"

    def _run_conversion():
        # HOME=/tmp is often required in Docker — LibreOffice needs a writable profile dir.
        # We use tempfile.gettempdir() to be more portable while staying secure.
        env = os.environ.copy()
        env["HOME"] = tempfile.gettempdir()
        
        # Find absolute path for libreoffice to satisfy Bandit B607
        libreoffice_path = shutil.which("libreoffice") or "libreoffice"
        
        result = subprocess.run(  # nosec B603
            [
                libreoffice_path,
                "--headless",
                "--norestore",
                "--nofirststartwizard",
                "--convert-to", "pdf",
                "--outdir", abs_outdir,
                abs_pptx,
            ],
            capture_output=True,
            timeout=120,
            env=env,
        )
        return result.returncode, result.stdout.decode(errors="replace"), result.stderr.decode(errors="replace")

    try:
        loop = asyncio.get_event_loop()
        returncode, stdout, stderr = await loop.run_in_executor(None, _run_conversion)

        logger.debug(f"LibreOffice stdout: {stdout}")
        if stderr:
            logger.debug(f"LibreOffice stderr: {stderr}")

        if returncode != 0:
            logger.warning(f"LibreOffice exited with code {returncode} for {abs_pptx}. stderr: {stderr}")
            return None

        # LibreOffice creates {abs_outdir}/{basename_without_ext}.pdf — rename to preview path
        basename_no_ext = os.path.splitext(os.path.basename(abs_pptx))[0]
        libreoffice_output = os.path.join(abs_outdir, basename_no_ext + ".pdf")

        if os.path.exists(libreoffice_output):
            os.rename(libreoffice_output, preview_path)
            # Strip bookmarks/outline so the browser navpane doesn't open
            _strip_pdf_bookmarks(preview_path)
            logger.info(f"PPTX preview PDF created: {preview_path}")
            return preview_path

        logger.warning(f"LibreOffice output not found at: {libreoffice_output}. stdout: {stdout}")
    except subprocess.TimeoutExpired:
        logger.warning(f"LibreOffice conversion timed out for {abs_pptx}")
    except FileNotFoundError:
        logger.warning("LibreOffice not found. Install it to enable PPTX preview.")
    except Exception as e:
        logger.warning(f"PPTX to PDF conversion failed for {abs_pptx}: {e}")

    return None


def get_pptx_orientation(file_path: str) -> tuple[str, float]:
    """
    Quickly detects the orientation of a PPTX file.
    """
    try:
        prs = Presentation(file_path)
        width = prs.slide_width
        height = prs.slide_height
        aspect_ratio = width / height if height > 0 else 1.777
        return "portrait" if height > width else "landscape", aspect_ratio
    except Exception as e:
        logger.warning(f"Failed to detect PPTX orientation: {e}")
    return "landscape", 1.777
